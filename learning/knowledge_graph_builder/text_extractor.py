"""
统一的文本提取模块
支持 PDF / DOCX / PPTX → 纯文本，供后续知识图谱三元组抽取使用
"""

import os
import sys
import io

_orig_stdout = sys.stdout
if hasattr(sys.stdout, 'buffer') and sys.stdout.buffer and not sys.stdout.buffer.closed:
    try:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    except Exception:
        sys.stdout = _orig_stdout

from pathlib import Path


def extract_text_from_file(file_path: str, file_type: str = '') -> str:
    """
    根据文件类型自动提取文本

    参数:
        file_path: 文件路径
        file_type: 文件类型标识（'pdf', 'docx', 'doc', 'pptx', 'ppt'），
                   为空时从扩展名自动判断

    返回:
        提取的纯文本内容

    抛出:
        ValueError: 不支持的文件类型
        FileNotFoundError: 文件不存在
        ImportError: 缺少依赖库
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    if not file_type:
        ext = path.suffix.lower()
        type_map = {
            '.pdf': 'pdf', '.docx': 'docx', '.doc': 'doc',
            '.pptx': 'pptx', '.ppt': 'ppt',
        }
        file_type = type_map.get(ext, '')

    if file_type in ('pdf',):
        return _extract_pdf(str(path))
    elif file_type in ('docx', 'doc'):
        return _extract_docx(str(path))
    elif file_type in ('pptx', 'ppt'):
        return _extract_pptx(str(path))
    else:
        raise ValueError(f"不支持的文本提取文件类型: {file_type}")


def _extract_pdf(pdf_path: str) -> str:
    """使用已有 pdf_extractor 模块提取 PDF 文本"""
    try:
        from learning.knowledge_graph_builder.pdf_extractor import extract_text_from_pdf
    except ImportError:
        # 直接调用（当从脚本运行时）
        try:
            from pdf_extractor import extract_text_from_pdf
        except ImportError:
            raise ImportError("无法导入 PDF 提取模块")

    return extract_text_from_pdf(pdf_path)


def _extract_docx(docx_path: str) -> str:
    """提取 DOCX 文档文本"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")

    doc = Document(docx_path)
    lines = []

    # 段落
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            lines.append(text)

    # 表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = ' | '.join(c for c in cells if c)
            if line:
                lines.append(line)

    return '\n'.join(lines)


def _extract_pptx(pptx_path: str) -> str:
    """提取 PPTX 文档文本"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    prs = Presentation(pptx_path)
    lines = []

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"--- 第 {slide_num} 页 ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = ' | '.join(c for c in cells if c)
                    if line:
                        lines.append(line)
        lines.append('')

    return '\n'.join(lines)
